import os
import json
from pathlib import Path
from datetime import datetime
from typing import Optional

try:
    from .models import AssessmentData
except ImportError:
    from models import AssessmentData


def _timestamp() -> str:
    return datetime.now().strftime("%Y-%m-%d_%H%M%S")


def _score_avg(items: list) -> Optional[float]:
    scored = [i for i in items if i.get("score") is not None and not i.get("na", False)]
    if not scored:
        return None
    return sum(i["score"] for i in scored) / len(scored)


def _get_maturity_band(score: float) -> dict:
    bands = [
        {"min": 1.0, "max": 1.5, "label": "Reactive", "color": "#ef4444"},
        {"min": 1.5, "max": 2.0, "label": "Emerging", "color": "#f97316"},
        {"min": 2.0, "max": 2.5, "label": "Developing", "color": "#eab308"},
        {"min": 2.5, "max": 3.0, "label": "Established", "color": "#84cc16"},
        {"min": 3.0, "max": 3.5, "label": "Managed", "color": "#22c55e"},
        {"min": 3.5, "max": 4.0, "label": "Optimizing", "color": "#15803d"},
    ]
    for band in bands:
        if band["min"] <= score < band["max"]:
            return band
    if score >= 4.0:
        return bands[-1]
    return bands[0]


def _score_label(score: int) -> str:
    labels = {1: "Initial", 2: "Developing", 3: "Established", 4: "Optimizing"}
    return labels.get(score, str(score))


def _weighted_composite(data_dict: dict) -> Optional[float]:
    weights = data_dict.get("scoring_config", {}).get("pillar_weights", {})
    total_weight = 0
    weighted_sum = 0
    for pillar in data_dict.get("pillars", []):
        items = []
        for ca in pillar.get("capability_areas", []):
            items.extend(ca.get("items", []))
        score = _score_avg(items)
        weight = weights.get(pillar["id"], 0)
        if score is not None:
            weighted_sum += score * weight
            total_weight += weight
    if total_weight == 0:
        return None
    return weighted_sum / total_weight


# TOGAF ADM phase mapping for pillars
TOGAF_MAPPING = {
    "governance": {
        "phases": ["Preliminary Phase", "Phase A: Architecture Vision"],
        "requirement": "Architecture governance frameworks and compliance",
    },
    "strategy": {
        "phases": ["Phase A: Architecture Vision", "Phase B: Business Architecture"],
        "requirement": "Strategic alignment and business direction",
    },
    "enterprise-architecture": {
        "phases": ["Phase B: Business Architecture", "Phase C: Information Systems", "Phase D: Technology Architecture"],
        "requirement": "Enterprise architecture management and TOGAF ADM",
    },
    "service-management": {
        "phases": ["Phase E: Opportunities & Solutions", "Phase F: Migration Planning"],
        "requirement": "IT service delivery and ITIL alignment",
    },
    "infrastructure-cloud": {
        "phases": ["Phase D: Technology Architecture", "Phase E: Opportunities & Solutions"],
        "requirement": "Infrastructure modernization and cloud adoption",
    },
    "workforce-culture": {
        "phases": ["Phase B: Business Architecture", "Phase H: Architecture Change Management"],
        "requirement": "Workforce capability and organizational culture",
    },
    "innovation-digital": {
        "phases": ["Phase A: Architecture Vision", "Phase H: Architecture Change Management"],
        "requirement": "Digital transformation and innovation management",
    },
}


class ExportEngine:
    def __init__(self, base_dir: str, resource_dir: str | None = None):
        self.base_dir = Path(base_dir)
        self.exports_dir = self.base_dir / "exports"

    def _ensure_exports_dir(self):
        self.exports_dir.mkdir(exist_ok=True)

    def _data_dict(self, data: AssessmentData) -> dict:
        return json.loads(data.model_dump_json())

    def generate_radar_chart_png(self, data: AssessmentData) -> str:
        """Generate radar chart as PNG, return path."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        d = self._data_dict(data)
        labels = []
        scores = []
        for pillar in d.get("pillars", []):
            # Shorten long names for the chart
            name = pillar["name"]
            if len(name) > 18:
                # Wrap at a space near the midpoint
                mid = len(name) // 2
                space = name.rfind(" ", 0, mid + 6)
                if space > 0:
                    name = name[:space] + "\n" + name[space + 1:]
            labels.append(name)
            items = []
            for ca in pillar.get("capability_areas", []):
                items.extend(ca.get("items", []))
            s = _score_avg(items)
            scores.append(s if s is not None else 0)

        n = len(labels)
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        scores_plot = scores + [scores[0]]
        angles += [angles[0]]

        fig, ax = plt.subplots(figsize=(7, 7), subplot_kw=dict(polar=True))
        ax.fill(angles, scores_plot, alpha=0.25, color="#3b82f6")
        ax.plot(angles, scores_plot, color="#3b82f6", linewidth=2)
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(labels, size=9)
        ax.set_ylim(0, 4)
        ax.set_yticks([1, 2, 3, 4])
        ax.set_yticklabels(["1\nInitial", "2\nDeveloping", "3\nEstablished", "4\nOptimizing"], size=7)
        ax.set_title("IT Strategy Maturity Profile", size=14, pad=25)

        self._ensure_exports_dir()
        path = str(self.exports_dir / "radar_chart.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", transparent=False, facecolor="white")
        plt.close(fig)
        return path

    # ------------------------------------------------------------------
    # D-01: Assessment Findings
    # ------------------------------------------------------------------

    def export_findings(self, data: AssessmentData) -> str:
        """D-01: Assessment Findings (Word)."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-01_AssessmentFindings_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        try:
            self._create_findings_docx(d, str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-01 Assessment Findings export failed: {e}") from e

        return filename

    def _create_findings_docx(self, d: dict, output_path: str):
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        title = doc.add_heading("IT Strategy Assessment Findings", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Assessment Date: {info.get('assessment_date', '')}")
        doc.add_paragraph(f"Assessor: {info.get('assessor', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")

        composite = _weighted_composite(d)
        if composite:
            band = _get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            p = doc.add_paragraph()
            p.add_run(f"Composite Score: {composite:.2f}").bold = True
            doc.add_paragraph(f"Maturity Band: {band['label']}")

        doc.add_heading("Pillar Findings", level=1)

        for pillar in d.get("pillars", []):
            items_all = []
            for ca in pillar.get("capability_areas", []):
                items_all.extend(ca.get("items", []))
            score = _score_avg(items_all)

            doc.add_heading(pillar["name"], level=2)
            weight_pct = f"{pillar.get('weight', 0) * 100:.0f}%"
            if score:
                band = _get_maturity_band(score)
                doc.add_paragraph(f"Score: {score:.2f} ({band['label']})  |  Weight: {weight_pct}")
            else:
                doc.add_paragraph(f"Score: Not yet scored  |  Weight: {weight_pct}")

            for ca in pillar.get("capability_areas", []):
                ca_score = _score_avg(ca.get("items", []))
                doc.add_heading(ca["name"], level=3)
                if ca_score:
                    doc.add_paragraph(f"Capability Area Score: {ca_score:.2f} — {_get_maturity_band(ca_score)['label']}")
                else:
                    doc.add_paragraph("Capability Area Score: Not scored")

                for item in ca.get("items", []):
                    if item.get("na"):
                        label = f"[N/A] {item['text']}"
                        if item.get("na_justification"):
                            label += f"\n  Justification: {item['na_justification']}"
                    else:
                        score_val = item.get("score")
                        score_str = f"{score_val} — {_score_label(score_val)}" if score_val else "Not scored"
                        label = f"[{score_str}] {item['text']}"
                    if item.get("notes"):
                        label += f"\n  Notes: {item['notes']}"
                    if item.get("confidence"):
                        label += f"  (Confidence: {item['confidence']})"
                    doc.add_paragraph(label, style="List Bullet")

        # Gov Compliance section
        if d.get("gov_compliance_enabled") and d.get("gov_compliance_extension"):
            gce = d["gov_compliance_extension"]
            doc.add_heading("Government Compliance Extension", level=1)

            for sub_key, sub_label in [("federal", "Federal Compliance"), ("state", "State Compliance")]:
                sub = gce.get(sub_key, {})
                if not sub.get("enabled"):
                    continue
                doc.add_heading(sub_label, level=2)
                for section in sub.get("sections", []):
                    doc.add_heading(section["name"], level=3)
                    for ca in section.get("capability_areas", []):
                        ca_score = _score_avg(ca.get("items", []))
                        doc.add_heading(ca["name"], level=4)
                        if ca_score:
                            doc.add_paragraph(f"Score: {ca_score:.2f} — {_get_maturity_band(ca_score)['label']}")
                        for item in ca.get("items", []):
                            score_val = item.get("score")
                            score_str = f"{score_val} — {_score_label(score_val)}" if score_val and not item.get("na") else ("N/A" if item.get("na") else "Not scored")
                            doc.add_paragraph(f"[{score_str}] {item['text']}", style="List Bullet")

        doc.save(output_path)

    # ------------------------------------------------------------------
    # D-02: Executive Summary
    # ------------------------------------------------------------------

    def export_executive_summary(self, data: AssessmentData) -> str:
        """D-02: Executive Summary (Word)."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-02_ExecutiveSummary_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        chart_path = self.generate_radar_chart_png(data)

        try:
            from docx import Document
            from docx.shared import Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            title = doc.add_heading("IT Strategy Assessment — Executive Summary", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            info = d.get("client_info", {})
            doc.add_paragraph(f"Client: {info.get('name', '')}")
            doc.add_paragraph(f"Assessment Date: {info.get('assessment_date', '')}")
            doc.add_paragraph(f"Assessor: {info.get('assessor', '')}")

            composite = _weighted_composite(d)
            if composite:
                band = _get_maturity_band(composite)
                doc.add_heading("Overall Maturity", level=1)
                p = doc.add_paragraph()
                p.add_run(f"Composite Score: {composite:.2f}  |  Band: {band['label']}").bold = True

                doc.add_paragraph(
                    "This score reflects a weighted composite across all seven TOGAF-aligned pillars. "
                    f"The organization is currently in the '{band['label']}' band, indicating "
                    + _band_narrative(band["label"])
                )

            doc.add_heading("Maturity Profile", level=1)
            doc.add_picture(chart_path, width=Inches(5.5))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

            doc.add_heading("Pillar Scores", level=1)
            table = doc.add_table(rows=1, cols=4)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            hdr[0].text = "Pillar"
            hdr[1].text = "Weight"
            hdr[2].text = "Score"
            hdr[3].text = "Band"
            for hc in hdr:
                for run in hc.paragraphs[0].runs:
                    run.bold = True

            for pillar in d.get("pillars", []):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)
                row = table.add_row().cells
                row[0].text = pillar["name"]
                row[1].text = f"{pillar.get('weight', 0) * 100:.0f}%"
                row[2].text = f"{score:.2f}" if score else "N/A"
                row[3].text = _get_maturity_band(score)["label"] if score else "N/A"

            doc.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-02 Executive Summary export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # D-03: Gap Analysis
    # ------------------------------------------------------------------

    def export_gap_analysis(self, data: AssessmentData) -> str:
        """D-03: Gap Analysis & Roadmap (Word)."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-03_GapAnalysis_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        target_scores = d.get("target_scores", {})
        gap_data = []
        for pillar in d.get("pillars", []):
            items_all = []
            for ca in pillar.get("capability_areas", []):
                items_all.extend(ca.get("items", []))
            current = _score_avg(items_all)
            target = target_scores.get(pillar["id"], 3.0)
            gap = (target - current) if current is not None else None
            gap_data.append({
                "pillar": pillar["name"],
                "pillar_id": pillar["id"],
                "current": round(current, 2) if current is not None else None,
                "target": target,
                "gap": round(gap, 2) if gap is not None else None,
                "severity": "High" if gap and gap > 1.5 else "Medium" if gap and gap > 0.5 else "Low",
                "capability_areas": pillar.get("capability_areas", []),
            })

        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            title = doc.add_heading("IT Strategy Assessment — Gap Analysis & Roadmap", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            info = d.get("client_info", {})
            doc.add_paragraph(f"Client: {info.get('name', '')}")
            doc.add_paragraph(f"Assessment Date: {info.get('assessment_date', '')}")

            doc.add_heading("Gap Summary", level=1)
            table = doc.add_table(rows=1, cols=5)
            table.style = "Table Grid"
            headers = ["Pillar", "Current Score", "Target Score", "Gap", "Severity"]
            for i, h in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = h
                for run in cell.paragraphs[0].runs:
                    run.bold = True

            for gap in gap_data:
                row = table.add_row().cells
                row[0].text = gap["pillar"]
                row[1].text = str(gap["current"]) if gap["current"] is not None else "N/A"
                row[2].text = str(gap["target"])
                row[3].text = str(gap["gap"]) if gap["gap"] is not None else "N/A"
                row[4].text = gap["severity"] if gap["gap"] is not None else "N/A"

            doc.add_heading("Detailed Gap Analysis by Pillar", level=1)

            for gap in gap_data:
                doc.add_heading(gap["pillar"], level=2)
                if gap["current"] is not None:
                    doc.add_paragraph(
                        f"Current: {gap['current']}  |  Target: {gap['target']}  |  "
                        f"Gap: {gap['gap']}  |  Severity: {gap['severity']}"
                    )
                else:
                    doc.add_paragraph("Not yet scored.")

                doc.add_heading("Capability Area Gaps", level=3)
                for ca in gap["capability_areas"]:
                    ca_score = _score_avg(ca.get("items", []))
                    ca_gap = (gap["target"] - ca_score) if ca_score is not None else None
                    if ca_score is not None:
                        doc.add_paragraph(
                            f"{ca['name']}: {ca_score:.2f} → target {gap['target']:.1f} "
                            f"(gap: {ca_gap:.2f})",
                            style="List Bullet",
                        )
                    else:
                        doc.add_paragraph(f"{ca['name']}: Not scored", style="List Bullet")

                # Low-scoring items as recommendations
                low_items = []
                for ca in gap["capability_areas"]:
                    for item in ca.get("items", []):
                        s = item.get("score")
                        if s is not None and s <= 2 and not item.get("na"):
                            low_items.append((s, item["text"], ca["name"]))
                if low_items:
                    low_items.sort(key=lambda x: x[0])
                    doc.add_heading("Recommended Focus Areas", level=3)
                    for s, text, ca_name in low_items[:5]:
                        doc.add_paragraph(
                            f"[Score {s} — {_score_label(s)}] {text} ({ca_name})",
                            style="List Bullet",
                        )

            doc.add_heading("Roadmap Priorities", level=1)
            high_gaps = [g for g in gap_data if g["gap"] and g["gap"] > 0.5]
            high_gaps.sort(key=lambda x: -(x["gap"] or 0))
            if high_gaps:
                for g in high_gaps:
                    doc.add_paragraph(
                        f"{g['pillar']}: Close gap of {g['gap']} points to reach target {g['target']}",
                        style="List Number",
                    )
            else:
                doc.add_paragraph("No significant gaps identified. All pillars are at or near target scores.")

            doc.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-03 Gap Analysis export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # D-04: Scored Assessment Workbook
    # ------------------------------------------------------------------

    def export_workbook(self, data: AssessmentData) -> str:
        """D-04: Scored Assessment Workbook (Excel)."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-04_ScoredWorkbook_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        try:
            self._create_workbook(d, str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-04 Scored Workbook export failed: {e}") from e

        return filename

    def _create_workbook(self, d: dict, output_path: str):
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        HEADER_FILL = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
        HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
        SUBHEADER_FILL = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")
        SUBHEADER_FONT = Font(bold=True, size=10)

        SCORE_COLORS = {
            1: "FEE2E2",  # red-100
            2: "FEF3C7",  # amber-100
            3: "D1FAE5",  # green-100
            4: "DBEAFE",  # blue-100
        }

        wb = Workbook()

        # ---- Dashboard sheet ----
        ws = wb.active
        ws.title = "Dashboard"

        ws.merge_cells("A1:F1")
        ws["A1"] = "IT Strategy Assessment Dashboard"
        ws["A1"].font = Font(bold=True, size=16, color="1E3A5F")
        ws["A1"].alignment = Alignment(horizontal="center")

        ws["A3"] = "Client:"
        ws["B3"] = d.get("client_info", {}).get("name", "")
        ws["A4"] = "Assessment Date:"
        ws["B4"] = d.get("client_info", {}).get("assessment_date", "")
        ws["A5"] = "Assessor:"
        ws["B5"] = d.get("client_info", {}).get("assessor", "")
        ws["A6"] = "Industry:"
        ws["B6"] = d.get("client_info", {}).get("industry", "")

        composite = _weighted_composite(d)
        ws["A8"] = "Overall Maturity Score:"
        ws["A8"].font = Font(bold=True)
        ws["B8"] = round(composite, 2) if composite is not None else "N/A"
        if composite is not None:
            band = _get_maturity_band(composite)
            ws["C8"] = band["label"]

        row = 10
        for cell_ref, val in [(f"A{row}", "Pillar"), (f"B{row}", "Weight"), (f"C{row}", "Score"), (f"D{row}", "Band"), (f"E{row}", "Target"), (f"F{row}", "Gap")]:
            ws[cell_ref] = val
            ws[cell_ref].font = HEADER_FONT
            ws[cell_ref].fill = HEADER_FILL

        target_scores = d.get("target_scores", {})
        for pillar in d.get("pillars", []):
            row += 1
            items_all = []
            for ca in pillar.get("capability_areas", []):
                items_all.extend(ca.get("items", []))
            score = _score_avg(items_all)
            target = target_scores.get(pillar["id"])
            gap = round(target - score, 2) if (score is not None and target is not None) else None

            ws[f"A{row}"] = pillar["name"]
            ws[f"B{row}"] = f"{pillar.get('weight', 0) * 100:.0f}%"
            ws[f"C{row}"] = round(score, 2) if score is not None else "N/A"
            if score is not None:
                ws[f"D{row}"] = _get_maturity_band(score)["label"]
            ws[f"E{row}"] = target if target is not None else ""
            ws[f"F{row}"] = gap if gap is not None else ""

        ws.column_dimensions["A"].width = 30
        ws.column_dimensions["B"].width = 10
        ws.column_dimensions["C"].width = 10
        ws.column_dimensions["D"].width = 16
        ws.column_dimensions["E"].width = 10
        ws.column_dimensions["F"].width = 10

        # ---- Per-pillar sheets ----
        for pillar in d.get("pillars", []):
            ws = wb.create_sheet(title=pillar["name"][:31])
            ws["A1"] = pillar["name"]
            ws["A1"].font = Font(bold=True, size=13, color="1E3A5F")

            row = 3
            headers = ["Capability Area", "Item ID", "Assessment Item", "Score", "Label", "Confidence", "Notes", "Evidence"]
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=row, column=col, value=h)
                cell.font = HEADER_FONT
                cell.fill = HEADER_FILL

            for ca in pillar.get("capability_areas", []):
                for item in ca.get("items", []):
                    row += 1
                    score_val = item.get("score")
                    ws.cell(row=row, column=1, value=ca["name"])
                    ws.cell(row=row, column=2, value=item["id"])
                    ws.cell(row=row, column=3, value=item["text"])
                    ws.cell(row=row, column=4, value=score_val)
                    ws.cell(row=row, column=5, value=_score_label(score_val) if score_val else ("N/A" if item.get("na") else ""))
                    ws.cell(row=row, column=6, value=item.get("confidence", ""))
                    ws.cell(row=row, column=7, value=item.get("notes", ""))
                    refs = item.get("evidence_references", [])
                    evidence_str = "; ".join(
                        f"{r.get('document', '')} §{r.get('section', '')}"
                        for r in refs if r.get("document")
                    ) if refs else ""
                    ws.cell(row=row, column=8, value=evidence_str)

                    # Color by score
                    if score_val and score_val in SCORE_COLORS:
                        fill = PatternFill(start_color=SCORE_COLORS[score_val], end_color=SCORE_COLORS[score_val], fill_type="solid")
                        ws.cell(row=row, column=4).fill = fill

            # Auto-width
            col_widths = [30, 14, 60, 8, 14, 12, 40, 40]
            for i, w in enumerate(col_widths, 1):
                from openpyxl.utils import get_column_letter
                ws.column_dimensions[get_column_letter(i)].width = w

        # ---- Gov Compliance sheet (if enabled) ----
        if d.get("gov_compliance_enabled") and d.get("gov_compliance_extension"):
            gce = d["gov_compliance_extension"]
            for sub_key, sub_label in [("federal", "Federal Compliance"), ("state", "State Compliance")]:
                sub = gce.get(sub_key, {})
                if not sub.get("enabled"):
                    continue
                ws = wb.create_sheet(title=sub_label[:31])
                ws["A1"] = sub_label
                ws["A1"].font = Font(bold=True, size=13, color="1E3A5F")
                row = 3
                headers = ["Section", "Capability Area", "Item ID", "Assessment Item", "Score", "Label", "Notes"]
                for col, h in enumerate(headers, 1):
                    cell = ws.cell(row=row, column=col, value=h)
                    cell.font = HEADER_FONT
                    cell.fill = HEADER_FILL
                for section in sub.get("sections", []):
                    for ca in section.get("capability_areas", []):
                        for item in ca.get("items", []):
                            row += 1
                            score_val = item.get("score")
                            ws.cell(row=row, column=1, value=section["name"])
                            ws.cell(row=row, column=2, value=ca["name"])
                            ws.cell(row=row, column=3, value=item["id"])
                            ws.cell(row=row, column=4, value=item["text"])
                            ws.cell(row=row, column=5, value=score_val)
                            ws.cell(row=row, column=6, value=_score_label(score_val) if score_val else ("N/A" if item.get("na") else ""))
                            ws.cell(row=row, column=7, value=item.get("notes", ""))

        wb.save(output_path)

    # ------------------------------------------------------------------
    # D-05: Out-Brief Presentation
    # ------------------------------------------------------------------

    def export_outbrief(self, data: AssessmentData) -> str:
        """D-05: Out-Brief Presentation (PowerPoint)."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-05_OutBrief_{_timestamp()}.pptx"
        output_path = self.exports_dir / filename

        chart_path = self.generate_radar_chart_png(data)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def add_text_box(slide, text, left, top, width, height, bold=False, size=18, color=None):
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = RGBColor(*color)
                return txBox

            # Slide 1: Title
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)
            # Background rectangle
            bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(2.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(30, 58, 95)
            bg.line.fill.background()
            add_text_box(slide, "IT Strategy Assessment", 0.5, 0.3, 12, 1, bold=True, size=36, color=(255, 255, 255))
            add_text_box(slide, "Out-Brief Presentation", 0.5, 1.2, 12, 0.8, size=22, color=(180, 210, 240))
            add_text_box(slide, d.get("client_info", {}).get("name", ""), 0.5, 2.8, 8, 0.7, bold=True, size=20)
            add_text_box(slide, f"Assessment Date: {d.get('client_info', {}).get('assessment_date', '')}", 0.5, 3.6, 8, 0.6, size=14)
            add_text_box(slide, f"Assessor: {d.get('client_info', {}).get('assessor', '')}", 0.5, 4.2, 8, 0.6, size=14)

            # Slide 2: Assessment Overview
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "Assessment Overview", 0.5, 0.2, 12, 0.8, bold=True, size=24, color=(30, 58, 95))

            composite = _weighted_composite(d)
            overview_lines = [
                f"Client: {d.get('client_info', {}).get('name', '')}",
                f"Date: {d.get('client_info', {}).get('assessment_date', '')}",
                f"Assessor: {d.get('client_info', {}).get('assessor', '')}",
                "",
            ]
            if composite:
                band = _get_maturity_band(composite)
                overview_lines.append(f"Overall Composite Score: {composite:.2f}")
                overview_lines.append(f"Maturity Band: {band['label']}")

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(overview_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

            # Pillar quick table on the right
            pillar_lines = []
            for pillar in d.get("pillars", []):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)
                pillar_lines.append(f"{pillar['name']}: {score:.2f} — {_get_maturity_band(score)['label']}" if score else f"{pillar['name']}: N/A")

            txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(6), Inches(5.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, line in enumerate(pillar_lines):
                if i == 0:
                    tf2.paragraphs[0].text = line
                else:
                    tf2.add_paragraph().text = line

            # Slide 3: Radar Chart
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "IT Strategy Maturity Profile", 0.5, 0.2, 12, 0.8, bold=True, size=24, color=(30, 58, 95))
            slide.shapes.add_picture(chart_path, Inches(2.5), Inches(0.8), Inches(8), Inches(6.5))

            # Per-pillar slides
            for pillar in d.get("pillars", []):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)

                slide = prs.slides.add_slide(blank_layout)
                add_text_box(slide, pillar["name"], 0.5, 0.2, 12, 0.7, bold=True, size=22, color=(30, 58, 95))
                if score:
                    band = _get_maturity_band(score)
                    add_text_box(slide, f"Score: {score:.2f} — {band['label']}  |  Weight: {pillar.get('weight', 0) * 100:.0f}%", 0.5, 0.9, 12, 0.5, size=14)
                else:
                    add_text_box(slide, "Not yet scored", 0.5, 0.9, 12, 0.5, size=14)

                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                first = True
                for ca in pillar.get("capability_areas", []):
                    ca_score = _score_avg(ca.get("items", []))
                    line = f"• {ca['name']}: {ca_score:.2f} — {_get_maturity_band(ca_score)['label']}" if ca_score else f"• {ca['name']}: N/A"
                    if first:
                        tf.paragraphs[0].text = line
                        first = False
                    else:
                        tf.add_paragraph().text = line

            # Slide: Recommendations Summary
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "Key Recommendations", 0.5, 0.2, 12, 0.7, bold=True, size=24, color=(30, 58, 95))
            target_scores = d.get("target_scores", {})
            gap_items = []
            for pillar in d.get("pillars", []):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)
                target = target_scores.get(pillar["id"], 3.0)
                if score is not None and score < target:
                    gap_items.append((target - score, pillar["name"]))
            gap_items.sort(reverse=True)

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            first = True
            for gap, name in gap_items[:7]:
                line = f"• Improve {name} (gap: {gap:.2f})"
                if first:
                    tf.paragraphs[0].text = line
                    first = False
                else:
                    tf.add_paragraph().text = line
            if not gap_items:
                tf.paragraphs[0].text = "All pillars meet or exceed target scores."

            prs.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-05 Out-Brief export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # D-06: Maturity Heatmap
    # ------------------------------------------------------------------

    def export_heatmap(self, data: AssessmentData) -> str:
        """D-06: Maturity Heatmap (Excel). Pillar × CA color-coded score grid."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-06_MaturityHeatmap_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

            SCORE_FILLS = {
                1: PatternFill(start_color="FCA5A5", end_color="FCA5A5", fill_type="solid"),   # red-300
                2: PatternFill(start_color="FCD34D", end_color="FCD34D", fill_type="solid"),   # amber-300
                3: PatternFill(start_color="6EE7B7", end_color="6EE7B7", fill_type="solid"),   # emerald-300
                4: PatternFill(start_color="93C5FD", end_color="93C5FD", fill_type="solid"),   # blue-300
            }
            NA_FILL = PatternFill(start_color="E5E7EB", end_color="E5E7EB", fill_type="solid")
            HEADER_FILL = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
            HEADER_FONT = Font(bold=True, color="FFFFFF", size=10)

            thin = Side(border_style="thin", color="CBD5E1")
            border = Border(left=thin, right=thin, top=thin, bottom=thin)

            wb = Workbook()
            ws = wb.active
            ws.title = "Heatmap"

            # Title
            ws.merge_cells("A1:H1")
            ws["A1"] = "IT Strategy Assessment — Maturity Heatmap"
            ws["A1"].font = Font(bold=True, size=14, color="1E3A5F")
            ws["A1"].alignment = Alignment(horizontal="center")

            info = d.get("client_info", {})
            ws["A2"] = f"Client: {info.get('name', '')}  |  Date: {info.get('assessment_date', '')}"
            ws["A2"].font = Font(size=10, color="64748B")

            # Collect all unique CA names
            all_ca_names = []
            ca_names_set = set()
            for pillar in d.get("pillars", []):
                for ca in pillar.get("capability_areas", []):
                    if ca["name"] not in ca_names_set:
                        all_ca_names.append(ca["name"])
                        ca_names_set.add(ca["name"])

            pillars = d.get("pillars", [])

            # Header row: pillar names
            header_row = 4
            ws.cell(row=header_row, column=1, value="Capability Area").font = HEADER_FONT
            ws.cell(row=header_row, column=1).fill = HEADER_FILL

            for col_idx, pillar in enumerate(pillars, 2):
                cell = ws.cell(row=header_row, column=col_idx, value=pillar["name"])
                cell.font = HEADER_FONT
                cell.fill = HEADER_FILL
                cell.alignment = Alignment(horizontal="center", wrap_text=True)
                cell.border = border

            # Data rows: one per CA that appears in each pillar
            # We map by position: each pillar's CAs in order
            max_ca_count = max((len(p.get("capability_areas", [])) for p in pillars), default=0)

            for ca_row_idx in range(max_ca_count):
                row_num = header_row + 1 + ca_row_idx
                # CA name from first pillar that has this index
                first_ca_name = ""
                for pillar in pillars:
                    cas = pillar.get("capability_areas", [])
                    if ca_row_idx < len(cas):
                        first_ca_name = cas[ca_row_idx]["name"]
                        break
                ws.cell(row=row_num, column=1, value=first_ca_name).font = Font(size=9)
                ws.cell(row=row_num, column=1).border = border

                for col_idx, pillar in enumerate(pillars, 2):
                    cas = pillar.get("capability_areas", [])
                    cell = ws.cell(row=row_num, column=col_idx)
                    cell.border = border
                    cell.alignment = Alignment(horizontal="center")
                    if ca_row_idx < len(cas):
                        ca = cas[ca_row_idx]
                        ca_score = _score_avg(ca.get("items", []))
                        if ca_score is not None:
                            score_rounded = round(ca_score)
                            score_rounded = max(1, min(4, score_rounded))
                            cell.value = round(ca_score, 2)
                            cell.fill = SCORE_FILLS[score_rounded]
                            cell.font = Font(size=9, bold=True)
                        else:
                            cell.value = "—"
                            cell.fill = NA_FILL
                            cell.font = Font(size=9, color="9CA3AF")
                    else:
                        cell.value = ""
                        cell.fill = NA_FILL

            # Pillar average row
            avg_row = header_row + 1 + max_ca_count + 1
            ws.cell(row=avg_row, column=1, value="Pillar Average").font = Font(bold=True, size=9)
            ws.cell(row=avg_row, column=1).border = border
            for col_idx, pillar in enumerate(pillars, 2):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)
                cell = ws.cell(row=avg_row, column=col_idx)
                cell.border = border
                cell.alignment = Alignment(horizontal="center")
                if score is not None:
                    score_rounded = max(1, min(4, round(score)))
                    cell.value = round(score, 2)
                    cell.fill = SCORE_FILLS[score_rounded]
                    cell.font = Font(bold=True, size=9)
                else:
                    cell.value = "N/A"
                    cell.fill = NA_FILL

            # Legend
            legend_row = avg_row + 3
            ws.cell(row=legend_row, column=1, value="Legend:").font = Font(bold=True, size=9)
            for i, (score, label, color) in enumerate([
                (1, "Initial (1)", "FCA5A5"),
                (2, "Developing (2)", "FCD34D"),
                (3, "Established (3)", "6EE7B7"),
                (4, "Optimizing (4)", "93C5FD"),
            ], 1):
                cell = ws.cell(row=legend_row, column=i + 1, value=label)
                cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
                cell.font = Font(size=9, bold=True)
                cell.alignment = Alignment(horizontal="center")
                cell.border = border

            # Column widths
            ws.column_dimensions["A"].width = 35
            for col_idx in range(2, len(pillars) + 2):
                from openpyxl.utils import get_column_letter
                ws.column_dimensions[get_column_letter(col_idx)].width = 18
            ws.row_dimensions[header_row].height = 45

            wb.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-06 Maturity Heatmap export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # D-07: Quick Wins Report
    # ------------------------------------------------------------------

    def export_quick_wins(self, data: AssessmentData) -> str:
        """D-07: Quick Wins Report (Word). Items scored 1 or 2, grouped by Immediate/Short-term."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-07_QuickWins_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            # Collect all items scored 1 or 2
            quick_items = []
            for pillar in d.get("pillars", []):
                for ca in pillar.get("capability_areas", []):
                    for item in ca.get("items", []):
                        score = item.get("score")
                        if score in (1, 2) and not item.get("na"):
                            quick_items.append({
                                "score": score,
                                "pillar": pillar["name"],
                                "pillar_weight": pillar.get("weight", 0),
                                "ca": ca["name"],
                                "text": item["text"],
                                "id": item["id"],
                                "notes": item.get("notes", ""),
                                "confidence": item.get("confidence", ""),
                            })

            # Sort: score ascending (1 first), then pillar weight descending
            quick_items.sort(key=lambda x: (x["score"], -x["pillar_weight"]))

            immediate = [i for i in quick_items if i["score"] == 1]
            short_term = [i for i in quick_items if i["score"] == 2]

            doc = Document()
            title = doc.add_heading("IT Strategy Assessment — Quick Wins Report", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            info = d.get("client_info", {})
            doc.add_paragraph(f"Client: {info.get('name', '')}")
            doc.add_paragraph(f"Assessment Date: {info.get('assessment_date', '')}")

            doc.add_paragraph(
                "This report identifies assessment items scored at Initial (1) or Developing (2), "
                "representing opportunities for near-term improvement. Immediate actions address the most "
                "critical gaps; Short-term actions can be tackled within 3–6 months."
            )

            total = len(quick_items)
            doc.add_paragraph(
                f"Total quick-win opportunities identified: {total} "
                f"({len(immediate)} Immediate, {len(short_term)} Short-term)"
            )

            def write_section(items, heading_text, score_label_str):
                doc.add_heading(heading_text, level=1)
                if not items:
                    doc.add_paragraph(f"No {score_label_str} items identified.")
                    return
                doc.add_paragraph(
                    f"{len(items)} items scored at {score_label_str} level. "
                    "Sorted by pillar strategic weight (highest impact first)."
                )
                for item in items:
                    doc.add_heading(f"{item['id']}: {item['text'][:80]}{'...' if len(item['text']) > 80 else ''}", level=2)
                    table = doc.add_table(rows=4, cols=2)
                    table.style = "Table Grid"
                    rows_data = [
                        ("Pillar", item["pillar"]),
                        ("Capability Area", item["ca"]),
                        ("Current Score", f"{item['score']} — {_score_label(item['score'])}"),
                        ("Confidence", item["confidence"] or "Not specified"),
                    ]
                    for r_idx, (label, val) in enumerate(rows_data):
                        table.rows[r_idx].cells[0].text = label
                        table.rows[r_idx].cells[1].text = val
                        for run in table.rows[r_idx].cells[0].paragraphs[0].runs:
                            run.bold = True
                    if item["notes"]:
                        doc.add_paragraph(f"Notes: {item['notes']}")
                    doc.add_paragraph("")

            write_section(immediate, "Immediate Actions (Score = 1: Initial)", "Initial (1)")
            write_section(short_term, "Short-term Actions (Score = 2: Developing)", "Developing (2)")

            if not quick_items:
                doc.add_heading("Assessment Status", level=1)
                doc.add_paragraph(
                    "No items scored at Initial or Developing level were found. "
                    "The organization has achieved Established or higher across all scored items."
                )

            doc.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-07 Quick Wins export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # D-08: Compliance Mapping
    # ------------------------------------------------------------------

    def export_compliance_mapping(self, data: AssessmentData) -> str:
        """D-08: Compliance Mapping (Word). TOGAF ADM + optional FITARA/NASCIO mapping."""
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-08_ComplianceMapping_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            title = doc.add_heading("IT Strategy Assessment — Compliance Mapping", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            info = d.get("client_info", {})
            doc.add_paragraph(f"Client: {info.get('name', '')}")
            doc.add_paragraph(f"Assessment Date: {info.get('assessment_date', '')}")

            doc.add_paragraph(
                "This report maps the IT Strategy Assessment results to applicable compliance frameworks, "
                "standards, and best practices including TOGAF 10 ADM phases. "
                "Average scores and gaps are provided to prioritize remediation."
            )

            # ---- TOGAF ADM Mapping ----
            doc.add_heading("TOGAF 10 ADM Phase Mapping", level=1)
            doc.add_paragraph(
                "The following table maps each pillar's current maturity score to the relevant "
                "TOGAF Architecture Development Method (ADM) phases."
            )

            table = doc.add_table(rows=1, cols=5)
            table.style = "Table Grid"
            headers = ["Pillar / Requirement", "TOGAF ADM Phase(s)", "Avg Score", "Target", "Gap"]
            for i, h in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = h
                for run in cell.paragraphs[0].runs:
                    run.bold = True

            target_scores = d.get("target_scores", {})
            for pillar in d.get("pillars", []):
                items_all = []
                for ca in pillar.get("capability_areas", []):
                    items_all.extend(ca.get("items", []))
                score = _score_avg(items_all)
                target = target_scores.get(pillar["id"], 3.0)
                gap = round(target - score, 2) if score is not None else None
                togaf = TOGAF_MAPPING.get(pillar["id"], {})
                phases = ", ".join(togaf.get("phases", []))

                row = table.add_row().cells
                row[0].text = pillar["name"]
                row[1].text = phases
                row[2].text = f"{score:.2f}" if score is not None else "N/A"
                row[3].text = str(target)
                row[4].text = str(gap) if gap is not None else "N/A"

            # ---- Gov Compliance Mapping ----
            if d.get("gov_compliance_enabled") and d.get("gov_compliance_extension"):
                gce = d["gov_compliance_extension"]

                federal = gce.get("federal", {})
                if federal.get("enabled"):
                    doc.add_heading("Federal Compliance Mapping (FITARA / OMB Policy)", level=1)
                    doc.add_paragraph(
                        "The following table maps Federal compliance assessment results to FITARA requirements "
                        "and relevant OMB policy directives."
                    )
                    table = doc.add_table(rows=1, cols=5)
                    table.style = "Table Grid"
                    for i, h in enumerate(["Requirement", "Related Items", "Avg Score", "Gap", "Recommendation"]):
                        cell = table.rows[0].cells[i]
                        cell.text = h
                        for run in cell.paragraphs[0].runs:
                            run.bold = True
                    for section in federal.get("sections", []):
                        for ca in section.get("capability_areas", []):
                            ca_score = _score_avg(ca.get("items", []))
                            item_count = len(ca.get("items", []))
                            gap_val = round(3.0 - ca_score, 2) if ca_score is not None else None
                            recommendation = _compliance_recommendation(ca_score)
                            row = table.add_row().cells
                            row[0].text = f"{section['name']} — {ca['name']}"
                            row[1].text = str(item_count)
                            row[2].text = f"{ca_score:.2f}" if ca_score is not None else "N/A"
                            row[3].text = str(gap_val) if gap_val is not None else "N/A"
                            row[4].text = recommendation

                state = gce.get("state", {})
                if state.get("enabled"):
                    doc.add_heading("State Compliance Mapping (NASCIO / State IT Policy)", level=1)
                    doc.add_paragraph(
                        "The following table maps State compliance assessment results to NASCIO priorities "
                        "and applicable state IT governance requirements."
                    )
                    table = doc.add_table(rows=1, cols=5)
                    table.style = "Table Grid"
                    for i, h in enumerate(["Requirement", "Related Items", "Avg Score", "Gap", "Recommendation"]):
                        cell = table.rows[0].cells[i]
                        cell.text = h
                        for run in cell.paragraphs[0].runs:
                            run.bold = True
                    for section in state.get("sections", []):
                        for ca in section.get("capability_areas", []):
                            ca_score = _score_avg(ca.get("items", []))
                            item_count = len(ca.get("items", []))
                            gap_val = round(3.0 - ca_score, 2) if ca_score is not None else None
                            recommendation = _compliance_recommendation(ca_score)
                            row = table.add_row().cells
                            row[0].text = f"{section['name']} — {ca['name']}"
                            row[1].text = str(item_count)
                            row[2].text = f"{ca_score:.2f}" if ca_score is not None else "N/A"
                            row[3].text = str(gap_val) if gap_val is not None else "N/A"
                            row[4].text = recommendation

            doc.save(str(output_path))
        except Exception as e:
            if output_path.exists():
                output_path.unlink()
            raise RuntimeError(f"D-08 Compliance Mapping export failed: {e}") from e

        return filename

    # ------------------------------------------------------------------
    # Dispatch and export_all
    # ------------------------------------------------------------------

    def export(self, export_type: str, data: AssessmentData) -> str:
        EXPORT_MAP = {
            "findings": self.export_findings,
            "executive-summary": self.export_executive_summary,
            "gap-analysis": self.export_gap_analysis,
            "workbook": self.export_workbook,
            "outbrief": self.export_outbrief,
            "heatmap": self.export_heatmap,
            "quick-wins": self.export_quick_wins,
            "compliance-mapping": self.export_compliance_mapping,
        }
        fn = EXPORT_MAP.get(export_type)
        if fn is None:
            raise ValueError(f"Unknown export type: {export_type}")
        return fn(data)

    def export_all(self, data: AssessmentData) -> list[str]:
        return [
            self.export_findings(data),
            self.export_executive_summary(data),
            self.export_gap_analysis(data),
            self.export_workbook(data),
            self.export_outbrief(data),
            self.export_heatmap(data),
            self.export_quick_wins(data),
            self.export_compliance_mapping(data),
        ]


def _band_narrative(band_label: str) -> str:
    narratives = {
        "Reactive": "foundational IT strategy capabilities are largely absent or ad hoc.",
        "Emerging": "initial IT strategy practices are being established but remain inconsistent.",
        "Developing": "IT strategy practices are defined but not yet consistently applied across the organization.",
        "Established": "IT strategy practices are consistently applied and producing measurable results.",
        "Managed": "IT strategy is well-managed with quantitative targets and continuous improvement.",
        "Optimizing": "IT strategy is fully optimized, adaptive, and driving enterprise innovation.",
    }
    return narratives.get(band_label, "IT strategy maturity is progressing.")


def _compliance_recommendation(score: Optional[float]) -> str:
    if score is None:
        return "Assessment required"
    if score < 1.5:
        return "Critical — Immediate remediation required"
    if score < 2.0:
        return "High priority — Establish formal policies and procedures"
    if score < 2.5:
        return "Medium priority — Formalize and document existing practices"
    if score < 3.0:
        return "Medium priority — Improve consistency and measurement"
    if score < 3.5:
        return "Low priority — Focus on optimization and automation"
    return "Compliant — Maintain and continuously improve"
